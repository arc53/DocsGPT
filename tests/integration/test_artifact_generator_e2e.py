"""End-to-end ArtifactGeneratorTool: live Jupyter gateway + ephemeral Postgres + local storage.

Launches a real ``jupyter kernelgateway`` (no Docker — the credential helper
hangs on dev machines), wires the tool to the ephemeral pytest-postgresql DB and
a temp-dir ``LocalStorage``, and drives create/edit/rewrite through the full
path: a spec is rendered in the kernel by the FIXED renderer, the produced file
is persisted as an artifact version (spec = source of truth), and the rendered
bytes are re-opened with the real library to assert structure.

Skips gracefully when the gateway binary, websocket-client, or a renderer
library is unavailable.
"""

from __future__ import annotations

import io
import shutil
import socket
import subprocess
import time
import uuid

import pytest

requests = pytest.importorskip("requests")
pytest.importorskip("websocket")  # websocket-client
pytest.importorskip("pptx")
pytest.importorskip("docx")
pytest.importorskip("openpyxl")

from application.agents.tools.artifact_generator import ArtifactGeneratorTool  # noqa: E402
from application.sandbox.jupyter_gateway import JupyterKernelGatewaySandbox  # noqa: E402
from application.sandbox.manager import SandboxManager  # noqa: E402
from application.sandbox.sandbox_creator import SandboxCreator  # noqa: E402
from application.storage.db.repositories.artifacts import ArtifactsRepository  # noqa: E402
from application.storage.local import LocalStorage  # noqa: E402
from application.storage.storage_creator import StorageCreator  # noqa: E402

_GATEWAY_BIN = shutil.which("jupyter-kernelgateway") or shutil.which("jupyter")

pytestmark = [
    pytest.mark.integration,
    pytest.mark.skipif(
        _GATEWAY_BIN is None,
        reason="jupyter kernel gateway not installed (pip install jupyter-kernel-gateway)",
    ),
]


def _free_port() -> int:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
        sock.bind(("127.0.0.1", 0))
        return sock.getsockname()[1]


def _gateway_cmd(port: int) -> list:
    if _GATEWAY_BIN.endswith("jupyter-kernelgateway"):
        base = [_GATEWAY_BIN]
    else:
        base = [_GATEWAY_BIN, "kernelgateway"]
    return base + [
        "--KernelGatewayApp.ip=127.0.0.1",
        f"--KernelGatewayApp.port={port}",
        "--ZMQChannelsWebsocketConnection.limit_rate=False",
    ]


@pytest.fixture(scope="module")
def gateway_url():
    port = _free_port()
    proc = subprocess.Popen(_gateway_cmd(port), stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    url = f"http://127.0.0.1:{port}"
    deadline = time.time() + 30
    ready = False
    try:
        while time.time() < deadline:
            if proc.poll() is not None:
                pytest.skip("jupyter kernelgateway process exited during startup")
            try:
                if requests.get(f"{url}/api", timeout=1).status_code == 200:
                    ready = True
                    break
            except requests.RequestException:
                time.sleep(0.3)
        if not ready:
            pytest.skip("jupyter kernelgateway did not become ready in time")
        yield url
    finally:
        proc.terminate()
        try:
            proc.wait(timeout=10)
        except subprocess.TimeoutExpired:
            proc.kill()


@pytest.fixture()
def wired_tool(gateway_url, pg_engine, tmp_path, monkeypatch):
    """An ArtifactGeneratorTool wired to the live gateway, ephemeral PG, and a temp local store."""
    backend = JupyterKernelGatewaySandbox(gateway_url=gateway_url, default_timeout=60.0)
    SandboxCreator._instance = SandboxManager(backend=backend, max_ttl=1200.0)

    storage = LocalStorage(base_dir=str(tmp_path))
    monkeypatch.setattr(StorageCreator, "_instance", storage, raising=False)

    monkeypatch.setattr("application.storage.db.session.get_engine", lambda: pg_engine)

    conversation_id = str(uuid.uuid4())
    tool = ArtifactGeneratorTool(
        tool_config={"conversation_id": conversation_id, "tool_id": str(uuid.uuid4())},
        user_id="user-e2e",
    )
    try:
        yield tool, conversation_id, pg_engine, storage
    finally:
        SandboxCreator.reset()


def test_create_edit_rewrite_presentation_versions(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool
    from pptx import Presentation

    # create -> v1
    spec = {"title": "Q3", "slides": [{"title": "Intro", "bullets": ["one", "two"]}]}
    created = tool.execute_action("create_artifact", kind="presentation", title="Deck", spec=spec)
    assert created["status"] == "ok", created
    assert created["version"] == 1
    assert created["mime_type"].endswith("presentationml.presentation")
    assert created["filename"] == "Deck.pptx"
    artifact_id = created["artifact_id"]
    assert tool.get_artifact_id("create_artifact") == artifact_id

    # v1 spec is the source of truth; the rendered .pptx re-opens with one slide.
    with pg_engine.connect() as conn:
        repo = ArtifactsRepository(conn)
        artifact = repo.get_artifact_in_parent(artifact_id, conversation_id=conversation_id)
        v1 = repo.get_version(artifact_id, 1)
    assert artifact["current_version"] == 1
    assert v1["spec"] == spec
    assert v1["produced_by"]["spec_kind"] == "presentation"
    prs = Presentation(storage.get_file(v1["storage_path"]))
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text == "Intro"

    # edit (merge-patch) -> v2: add a second slide, keep title.
    patch = {"slides": [{"title": "Intro", "bullets": ["one", "two"]}, {"title": "Outro"}]}
    edited = tool.execute_action("edit_artifact", id=artifact_id, spec_patch=patch)
    assert edited["status"] == "ok", edited
    assert edited["version"] == 2
    assert edited["artifact_id"] == artifact_id

    with pg_engine.connect() as conn:
        repo = ArtifactsRepository(conn)
        artifact = repo.get_artifact_in_parent(artifact_id, conversation_id=conversation_id)
        v1_again = repo.get_version(artifact_id, 1)
        v2 = repo.get_version(artifact_id, 2)
    assert artifact["current_version"] == 2
    # Append-only: v1 unchanged and still present.
    assert v1_again["spec"] == spec
    assert v2["spec"]["title"] == "Q3"
    assert len(v2["spec"]["slides"]) == 2
    prs2 = Presentation(storage.get_file(v2["storage_path"]))
    assert len(prs2.slides) == 2

    # rewrite -> v3: wholesale new spec.
    new_spec = {"title": "Fresh", "slides": [{"title": "A"}, {"title": "B"}, {"title": "C"}]}
    rewritten = tool.execute_action("rewrite_artifact", id=artifact_id, spec=new_spec)
    assert rewritten["status"] == "ok", rewritten
    assert rewritten["version"] == 3

    with pg_engine.connect() as conn:
        repo = ArtifactsRepository(conn)
        artifact = repo.get_artifact_in_parent(artifact_id, conversation_id=conversation_id)
        versions = repo.list_versions(artifact_id)
        v3 = repo.get_version(artifact_id, 3)
    assert artifact["current_version"] == 3
    assert [v["version"] for v in versions] == [1, 2, 3]
    assert v3["spec"] == new_spec
    assert len(Presentation(storage.get_file(v3["storage_path"])).slides) == 3


def test_create_document_docx(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool
    from docx import Document

    spec = {"title": "Report", "sections": [{"heading": "Intro", "paragraphs": ["Hello.", "World."]}]}
    created = tool.execute_action("create_artifact", kind="document", spec=spec)
    assert created["status"] == "ok", created
    assert created["mime_type"].endswith("wordprocessingml.document")

    with pg_engine.connect() as conn:
        version = ArtifactsRepository(conn).get_version(created["artifact_id"], 1)
    doc = Document(storage.get_file(version["storage_path"]))
    texts = [p.text for p in doc.paragraphs]
    assert "Hello." in texts and "World." in texts


def test_create_spreadsheet_xlsx(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool
    from openpyxl import load_workbook

    spec = {"sheets": [{"name": "Data", "rows": [["a", "b"], [1, 2]]}]}
    created = tool.execute_action("create_artifact", kind="spreadsheet", spec=spec)
    assert created["status"] == "ok", created
    assert created["mime_type"].endswith("spreadsheetml.sheet")

    with pg_engine.connect() as conn:
        version = ArtifactsRepository(conn).get_version(created["artifact_id"], 1)
    wb = load_workbook(storage.get_file(version["storage_path"]))
    ws = wb["Data"]
    assert ws.cell(row=1, column=1).value == "a"
    assert ws.cell(row=2, column=2).value == 2


def test_create_edit_html_versions(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool

    # create -> v1: an html report rendered to report.html (kind html, text/html).
    spec = {
        "title": "Findings <hi>",
        "blocks": [
            {"type": "heading", "text": "Overview", "level": 2},
            {"type": "paragraph", "text": "Body <script>alert(1)</script>."},
            {"type": "list", "ordered": False, "items": ["one", "two"]},
            {"type": "table", "headers": ["k", "v"], "rows": [["a", "1"]]},
        ],
    }
    created = tool.execute_action("create_artifact", kind="html", title="Report", spec=spec)
    assert created["status"] == "ok", created
    assert created["version"] == 1
    assert created["mime_type"] == "text/html"
    assert created["filename"] == "report.html"
    assert created["ref"] == "A1"
    artifact_id = created["artifact_id"]

    with pg_engine.connect() as conn:
        repo = ArtifactsRepository(conn)
        artifact = repo.get_artifact_in_parent(artifact_id, conversation_id=conversation_id)
        v1 = repo.get_version(artifact_id, 1)
    assert artifact["kind"] == "html"
    assert artifact["current_version"] == 1
    assert v1["spec"] == spec
    assert v1["produced_by"]["spec_kind"] == "html"

    # The stored bytes are a valid HTML doc with escaped title/blocks (no live tag).
    html_doc = storage.get_file(v1["storage_path"]).read().decode("utf-8")
    assert html_doc.startswith("<!doctype html>")
    assert "Findings &lt;hi&gt;" in html_doc
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in html_doc
    assert "<script>alert(1)</script>" not in html_doc
    assert "<li>one</li>" in html_doc and "<td>a</td>" in html_doc

    # edit (merge-patch) -> v2: replace the title, keep v1 intact.
    edited = tool.execute_action("edit_artifact", id="A1", spec_patch={"title": "Updated"})
    assert edited["status"] == "ok", edited
    assert edited["version"] == 2
    assert edited["artifact_id"] == artifact_id
    assert edited["ref"] == "A1"

    with pg_engine.connect() as conn:
        repo = ArtifactsRepository(conn)
        artifact = repo.get_artifact_in_parent(artifact_id, conversation_id=conversation_id)
        v1_again = repo.get_version(artifact_id, 1)
        v2 = repo.get_version(artifact_id, 2)
    assert artifact["current_version"] == 2
    # Append-only: v1 unchanged; v2 carries the merged spec re-rendered.
    assert v1_again["spec"] == spec
    assert v2["spec"]["title"] == "Updated"
    assert len(v2["spec"]["blocks"]) == 4
    v2_doc = storage.get_file(v2["storage_path"]).read().decode("utf-8")
    assert "<h1>Updated</h1>" in v2_doc


def test_invalid_spec_creates_no_artifact(wired_tool):
    tool, conversation_id, pg_engine, _storage = wired_tool

    out = tool.execute_action("create_artifact", kind="presentation", spec={"title": "no slides"})
    assert out["status"] == "error"
    assert "invalid presentation spec" in out["error"]

    with pg_engine.connect() as conn:
        rows = ArtifactsRepository(conn).list_artifacts(conversation_id=conversation_id)
    assert rows == []


def test_edit_cross_tenant_denied(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool

    # An artifact owned by a DIFFERENT conversation must not be editable here.
    other_conversation = str(uuid.uuid4())
    foreign_id = _seed_presentation(pg_engine, storage, other_conversation)

    out = tool.execute_action("edit_artifact", id=foreign_id, spec_patch={"title": "x"})
    assert out["status"] == "error"
    assert "not found in this conversation/run" in out["error"]

    # The foreign artifact stayed at v1 (no version appended).
    with pg_engine.connect() as conn:
        artifact = ArtifactsRepository(conn).get_artifact(foreign_id)
    assert artifact["current_version"] == 1


def test_short_ref_numbering_and_edit_by_ref(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool
    from pptx import Presentation

    # Successive creates in one conversation get A1, A2 (1-based, by created order).
    first = tool.execute_action(
        "create_artifact", kind="presentation", title="One", spec={"slides": [{"title": "a"}]}
    )
    assert first["status"] == "ok", first
    assert first["ref"] == "A1"
    second = tool.execute_action(
        "create_artifact", kind="document", spec={"sections": [{"heading": "h", "paragraphs": ["p"]}]}
    )
    assert second["status"] == "ok", second
    assert second["ref"] == "A2"

    # edit_artifact(id="A1") resolves the ref -> appends v2 to the FIRST artifact only.
    edited = tool.execute_action("edit_artifact", id="A1", spec_patch={"slides": [{"title": "a"}, {"title": "b"}]})
    assert edited["status"] == "ok", edited
    assert edited["version"] == 2
    assert edited["artifact_id"] == first["artifact_id"]
    # The append keeps the artifact's position, so its ref is stable.
    assert edited["ref"] == "A1"

    with pg_engine.connect() as conn:
        repo = ArtifactsRepository(conn)
        artifact = repo.get_artifact_in_parent(first["artifact_id"], conversation_id=conversation_id)
        v1 = repo.get_version(first["artifact_id"], 1)
        v2 = repo.get_version(first["artifact_id"], 2)
    assert artifact["current_version"] == 2
    # Append-only: v1 preserved with one slide, v2 has two.
    assert len(v1["spec"]["slides"]) == 1
    assert len(v2["spec"]["slides"]) == 2
    assert len(Presentation(storage.get_file(v2["storage_path"])).slides) == 2

    # The uuid path still works alongside refs.
    by_uuid = tool.execute_action("edit_artifact", id=first["artifact_id"], spec_patch={"title": "X"})
    assert by_uuid["status"] == "ok"
    assert by_uuid["version"] == 3


def test_short_ref_is_conversation_scoped(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool

    # Seed an artifact in a DIFFERENT conversation; its A1 must not resolve here.
    other_conversation = str(uuid.uuid4())
    _seed_presentation(pg_engine, storage, other_conversation)

    # This conversation has no artifacts yet, so A1 resolves to nothing (no cross-parent leak).
    out = tool.execute_action("edit_artifact", id="A1", spec_patch={"title": "x"})
    assert out["status"] == "error"
    assert "not found in this conversation/run" in out["error"]


def test_short_ref_out_of_range_is_clean_error(wired_tool):
    tool, conversation_id, pg_engine, storage = wired_tool

    created = tool.execute_action("create_artifact", kind="presentation", spec={"slides": [{"title": "a"}]})
    assert created["ref"] == "A1"
    # A2 has no artifact behind it -> clean not-found, no crash.
    out = tool.execute_action("edit_artifact", id="A2", spec_patch={"title": "x"})
    assert out["status"] == "error"
    assert "not found in this conversation/run" in out["error"]


def _seed_presentation(pg_engine, storage, conversation_id) -> str:
    """Create a presentation artifact (row + version + bytes) in a given conversation; return its id."""
    spec = {"slides": [{"title": "seed"}]}
    data = b"%PDF-stub-bytes"  # content is irrelevant to the authz check
    with pg_engine.begin() as conn:
        repo = ArtifactsRepository(conn)
        artifact = repo.create_artifact(
            "user-e2e",
            "presentation",
            conversation_id=conversation_id,
            title="seed",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="seed.pptx",
            storage_path=None,
            size=len(data),
            sha256="0" * 64,
            spec=spec,
            produced_by={"tool": "artifact_generator", "action": "create_artifact", "spec_kind": "presentation"},
        )
        artifact_id = str(artifact["id"])
        from sqlalchemy import text

        storage_path = f"inputs/user-e2e/artifacts/{artifact_id}/v1/seed.pptx"
        storage.save_file(io.BytesIO(data), storage_path)
        conn.execute(
            text(
                "UPDATE artifact_versions SET storage_path = :p "
                "WHERE artifact_id = CAST(:aid AS uuid) AND version = 1"
            ),
            {"p": storage_path, "aid": artifact_id},
        )
    return artifact_id
