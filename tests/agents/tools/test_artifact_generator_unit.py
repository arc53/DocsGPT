"""Unit tests for ArtifactGeneratorTool: spec validation, merge-patch, and renderer injection safety.

No DB or sandbox: these exercise the pure logic (schema gate + RFC 7386 merge)
and the FIXED renderer programs directly (executed in-process against a temp dir)
to prove a spec value containing Python/quotes is rendered as literal text.
"""

from __future__ import annotations

import json
import os
import tempfile
from types import SimpleNamespace

import pytest

from application.agents.tools.artifact_generator import (
    _KIND_INFO,
    _RENDERERS,
    ArtifactGeneratorTool,
    merge_patch,
)

pytest.importorskip("pptx")
pytest.importorskip("docx")
pytest.importorskip("openpyxl")
pytest.importorskip("reportlab")


def _tool():
    return ArtifactGeneratorTool(
        tool_config={"conversation_id": "conv-1", "tool_id": "t-1"}, user_id="u-1"
    )


# ---------------------------------------------------------------------------
# Spec validation
# ---------------------------------------------------------------------------


def test_validate_accepts_minimal_presentation():
    assert _tool()._validate("presentation", {"slides": [{"title": "x"}]}) is None


def test_validate_rejects_missing_required_key():
    err = _tool()._validate("presentation", {"title": "no slides"})
    assert err["status"] == "error"
    assert "invalid presentation spec" in err["error"]


def test_validate_rejects_unknown_key():
    err = _tool()._validate("document", {"sections": [], "bogus": 1})
    assert err["status"] == "error"


# ---------------------------------------------------------------------------
# Reversion download filename
# ---------------------------------------------------------------------------


def _patch_reversion(monkeypatch, tool):
    """Stub render + persist so _reversion runs without a sandbox or DB; capture kwargs."""
    import application.agents.tools.artifact_generator as ag

    captured: dict = {}

    def _fake_append(**kwargs):
        captured.update(kwargs)
        return {"artifact_id": kwargs["artifact_id"], "version": 2, "id": "v2"}

    monkeypatch.setattr(ag, "append_artifact_version", _fake_append)
    monkeypatch.setattr(tool, "_render", lambda kind, spec: {"data": b"x"})
    return captured


def test_reversion_preserves_title_in_filename(monkeypatch):
    """An edited/rewritten version keeps the original artifact's download name."""
    tool = _tool()
    captured = _patch_reversion(monkeypatch, tool)

    result = tool._reversion(
        "art-1", "presentation", {"slides": [{"title": "x"}]}, "edit_artifact", "Q3 Deck"
    )

    assert result["status"] == "ok"
    assert captured["filename"] == "Q3 Deck.pptx"


def test_reversion_without_title_falls_back_to_generic(monkeypatch):
    """A missing title still yields a valid generic filename (never crashes)."""
    tool = _tool()
    captured = _patch_reversion(monkeypatch, tool)

    tool._reversion("art-1", "presentation", {"slides": [{"title": "x"}]}, "rewrite_artifact", None)

    assert captured["filename"] == "artifact.pptx"


def test_validate_rejects_wrong_type():
    err = _tool()._validate("spreadsheet", {"sheets": "not-a-list"})
    assert err["status"] == "error"


def test_validate_rejects_non_object_spec():
    err = _tool()._validate("pdf", "just a string")
    assert err["status"] == "error"
    assert "spec must be a JSON object" in err["error"]


def test_create_rejects_unknown_kind():
    out = _tool()._create(kind="hologram", spec={})
    assert out["status"] == "error"
    assert "unsupported kind" in out["error"]


# ---------------------------------------------------------------------------
# _render session lifecycle: cleans its scratch dir, leaves the shared session open
# ---------------------------------------------------------------------------


class _FakeRenderManager:
    """Records remove_path/close and returns fixed render bytes; no real sandbox."""

    def __init__(self):
        self.closed = []
        self.removed = []

    def open(self, session_id, ttl=None):
        return session_id

    def put_file(self, session_id, dest_path, data):
        pass

    def exec(self, session_id, code, timeout=None):
        return SimpleNamespace(ok=True)

    def get_file(self, session_id, path):
        return b"%PDF-1.4 rendered"

    def remove_path(self, session_id, path):
        self.removed.append((session_id, path))

    def close(self, session_id):
        self.closed.append(session_id)


def test_render_cleans_scratch_but_leaves_session_open(monkeypatch):
    """_render drops its per-render scratch dir but must NOT close the shared session."""
    manager = _FakeRenderManager()
    monkeypatch.setattr(
        "application.sandbox.sandbox_creator.SandboxCreator.get_manager", lambda: manager
    )

    out = _tool()._render("pdf", {"title": "t", "blocks": []})

    assert out == {"data": b"%PDF-1.4 rendered"}
    # The render owns only its scratch dir; the session is the shared conversation
    # session that code_executor(persist=True) keeps warm, so it is left for the
    # manager/conversation to reap.
    assert manager.closed == []
    assert len(manager.removed) == 1
    session_id, path = manager.removed[0]
    assert session_id == "conv-1"
    assert path.startswith("artifacts/")


# ---------------------------------------------------------------------------
# RFC 7386 JSON merge-patch
# ---------------------------------------------------------------------------


def test_merge_patch_adds_and_overwrites():
    assert merge_patch({"a": 1, "b": 2}, {"b": 3, "c": 4}) == {"a": 1, "b": 3, "c": 4}


def test_merge_patch_null_deletes_key():
    assert merge_patch({"a": 1, "b": 2}, {"b": None}) == {"a": 1}


def test_merge_patch_recurses_into_objects():
    assert merge_patch({"x": {"a": 1, "b": 2}}, {"x": {"b": None, "c": 3}}) == {"x": {"a": 1, "c": 3}}


def test_merge_patch_replaces_array_wholesale():
    # RFC 7386: arrays are replaced, not merged element-wise.
    assert merge_patch({"l": [1, 2, 3]}, {"l": [9]}) == {"l": [9]}


def test_merge_patch_non_object_patch_replaces_target():
    assert merge_patch({"a": 1}, "scalar") == "scalar"


def test_merge_patch_does_not_mutate_inputs():
    target = {"a": {"b": 1}}
    patch = {"a": {"c": 2}}
    merge_patch(target, patch)
    assert target == {"a": {"b": 1}}
    assert patch == {"a": {"c": 2}}


# ---------------------------------------------------------------------------
# Renderer injection safety
# ---------------------------------------------------------------------------


def _render_in_process(kind: str, spec: dict) -> str:
    """Execute the FIXED renderer program against a temp dir; return the output path."""
    workdir = tempfile.mkdtemp()
    spec_path = os.path.join(workdir, "spec.json")
    out_path = os.path.join(workdir, f"out.{_KIND_INFO[kind]['ext']}")
    with open(spec_path, "w") as handle:
        json.dump(spec, handle)
    program = _RENDERERS[kind].format(spec_path=spec_path, out_path=out_path)
    namespace: dict = {}
    exec(compile(program, "<renderer>", "exec"), namespace, namespace)  # noqa: S102
    return out_path


def test_renderer_does_not_execute_spec_code(capfd, tmp_path):
    # A spec whose values are Python source / shell payloads must be treated as
    # literal text. If the renderer string-interpolated the spec it would run
    # this; instead it json.loads the spec as data, so nothing executes.
    sentinel = tmp_path / "pwned.txt"
    payload = (
        f"'''__import__('os').system('echo PWNED > {sentinel}')'''; "
        "print('SHOULD_NOT_PRINT')"
    )
    spec = {
        "title": payload,
        "slides": [{"title": payload, "bullets": [payload], "notes": payload}],
    }
    out_path = _render_in_process("presentation", spec)

    captured = capfd.readouterr()
    assert "SHOULD_NOT_PRINT" not in captured.out
    assert "PWNED" not in captured.out
    assert not sentinel.exists()
    assert os.path.getsize(out_path) > 0


def test_renderer_keeps_injection_text_as_literal_content():
    from pptx import Presentation

    payload = "'''; import os; os.system('echo HACK'); x = '''"
    spec = {"slides": [{"title": payload, "bullets": [payload]}]}
    prs = Presentation(_render_in_process("presentation", spec))
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text == payload


def test_spreadsheet_renderer_neutralizes_formula_injection():
    # Spec content is model / prompt-injection controlled. A string cell starting
    # with a formula trigger must be stored as TEXT (data_type 's'), never as a
    # live formula; genuine numbers stay numeric.
    from openpyxl import load_workbook

    spec = {
        "sheets": [
            {
                "name": "S",
                "rows": [
                    ['=HYPERLINK("http://evil","x")', "+1+1", "@cmd", "-danger"],
                    ["safe text", -5, 3.14],
                ],
            }
        ]
    }
    wb = load_workbook(_render_in_process("spreadsheet", spec))
    ws = wb["S"]
    # None of the trigger cells are formulas; each is a quote-prefixed string.
    for coord in ("A1", "B1", "C1", "D1"):
        assert ws[coord].data_type == "s", coord
        assert str(ws[coord].value).startswith("'")
    # Benign text is unchanged; numbers stay numeric (no spurious quoting).
    assert ws["A2"].value == "safe text" and ws["A2"].data_type == "s"
    assert ws["B2"].value == -5 and ws["B2"].data_type == "n"
    assert ws["C2"].value == 3.14 and ws["C2"].data_type == "n"


def test_pdf_renderer_escapes_markup_and_does_not_execute():
    payload = "<b>not bold</b> & \"</para>\" '''os.system('x')'''"
    spec = {"title": payload, "blocks": [{"type": "paragraph", "text": payload}]}
    out_path = _render_in_process("pdf", spec)
    assert os.path.getsize(out_path) > 0


# ---------------------------------------------------------------------------
# html kind: schema, renderer injection safety, output escaping
# ---------------------------------------------------------------------------


def test_validate_accepts_html_spec():
    spec = {
        "title": "Report",
        "blocks": [
            {"type": "heading", "text": "Intro", "level": 2},
            {"type": "paragraph", "text": "Body."},
            {"type": "list", "ordered": True, "items": ["a", "b"]},
            {"type": "table", "headers": ["h"], "rows": [["c"]]},
            {"type": "code", "text": "print(1)"},
        ],
    }
    assert _tool()._validate("html", spec) is None


def test_validate_rejects_html_unknown_block_key():
    err = _tool()._validate("html", {"blocks": [{"type": "paragraph", "text": "x", "bogus": 1}]})
    assert err["status"] == "error"


def test_validate_rejects_html_bad_heading_level():
    err = _tool()._validate("html", {"blocks": [{"type": "heading", "text": "x", "level": 9}]})
    assert err["status"] == "error"


def test_html_renderer_does_not_execute_spec_code(capfd, tmp_path):
    # The renderer json.loads the spec as data; a value that is Python source
    # must never run. If the program string-interpolated the spec it would.
    sentinel = tmp_path / "pwned.txt"
    payload = (
        f"'''__import__('os').system('echo PWNED > {sentinel}')'''; "
        "print('SHOULD_NOT_PRINT')"
    )
    spec = {"title": payload, "blocks": [{"type": "paragraph", "text": payload}]}
    out_path = _render_in_process("html", spec)

    captured = capfd.readouterr()
    assert "SHOULD_NOT_PRINT" not in captured.out
    assert "PWNED" not in captured.out
    assert not sentinel.exists()
    assert os.path.getsize(out_path) > 0


def test_html_renderer_escapes_spec_markup():
    # Spec text carrying live markup must appear HTML-ESCAPED in the output, so
    # no raw <script>/<img onerror> tag from spec content reaches the document.
    script = "<script>alert(1)</script>"
    img = '<img src=x onerror="alert(2)">'
    spec = {
        "title": script,
        "blocks": [
            {"type": "paragraph", "text": img},
            {"type": "heading", "text": "<b>h</b>", "level": 1},
            {"type": "list", "items": ["<i>x</i>"]},
            {"type": "table", "headers": ["<u>th</u>"], "rows": [["<em>td</em>"]]},
            {"type": "code", "text": "</code><script>x</script>"},
        ],
    }
    with open(_render_in_process("html", spec), encoding="utf-8") as handle:
        html_doc = handle.read()

    # No live tag from spec content survives.
    assert "<script>alert(1)</script>" not in html_doc
    assert "<img src=x onerror=" not in html_doc
    assert "<b>h</b>" not in html_doc
    assert "<i>x</i>" not in html_doc
    assert "<u>th</u>" not in html_doc
    assert "<em>td</em>" not in html_doc
    # The escaped forms are present instead.
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in html_doc
    assert "&lt;img src=x onerror=" in html_doc
    # The fixed structural HTML the renderer emits is intact.
    assert "<!doctype html>" in html_doc
    assert "<p>" in html_doc and "</p>" in html_doc
    assert "<table>" in html_doc and "<pre><code>" in html_doc


# ---------------------------------------------------------------------------
# Tool metadata: spec shapes are surfaced to the model
# ---------------------------------------------------------------------------


def _collect_schema_keys(schema, keys):
    if isinstance(schema, dict):
        for prop in schema.get("properties", {}):
            keys.add(prop)
        for sub in schema.values():
            _collect_schema_keys(sub, keys)
    elif isinstance(schema, list):
        for sub in schema:
            _collect_schema_keys(sub, keys)


def test_spec_synopsis_covers_every_schema_key():
    """Every kind and property in _SCHEMAS must appear in the metadata synopsis.

    The synopsis is a hand-written mirror of _SCHEMAS; this guards drift when a
    kind or key is added without updating what the model is told.
    """
    from application.agents.tools.artifact_generator import _SCHEMAS, _SPEC_SYNOPSIS

    for kind, schema in _SCHEMAS.items():
        assert kind in _SPEC_SYNOPSIS
        keys: set = set()
        _collect_schema_keys(schema, keys)
        for key in keys:
            assert key in _SPEC_SYNOPSIS, f"schema key {key!r} of kind {kind!r} missing from synopsis"


def test_create_and_rewrite_metadata_embed_spec_synopsis():
    from application.agents.tools.artifact_generator import _SPEC_SYNOPSIS

    actions = {a["name"]: a for a in _tool().get_actions_metadata()}
    create_spec = actions["create_artifact"]["parameters"]["properties"]["spec"]
    rewrite_spec = actions["rewrite_artifact"]["parameters"]["properties"]["spec"]
    assert _SPEC_SYNOPSIS in create_spec["description"]
    assert _SPEC_SYNOPSIS in rewrite_spec["description"]


# ---------------------------------------------------------------------------
# spec_append: additive edits that don't clobber arrays
# ---------------------------------------------------------------------------


def test_spec_append_preserves_existing_items():
    from application.agents.tools.artifact_generator import _apply_spec_append

    spec = {"title": "Brief", "blocks": [{"type": "heading", "text": "Overview"}]}
    out = _apply_spec_append(spec, {"blocks": [{"type": "heading", "text": "Risks"}]})
    assert "error" not in out
    assert [b["text"] for b in out["spec"]["blocks"]] == ["Overview", "Risks"]
    # The input spec is not mutated.
    assert len(spec["blocks"]) == 1


def test_spec_append_creates_missing_list():
    from application.agents.tools.artifact_generator import _apply_spec_append

    out = _apply_spec_append({"title": "x"}, {"blocks": [{"type": "paragraph", "text": "p"}]})
    assert out["spec"]["blocks"] == [{"type": "paragraph", "text": "p"}]


def test_spec_append_rejects_non_list_values_and_targets():
    from application.agents.tools.artifact_generator import _apply_spec_append

    assert "error" in _apply_spec_append({}, {"blocks": "not-a-list"})
    assert "error" in _apply_spec_append({"title": "t"}, {"title": ["x"]})


def test_edit_with_only_spec_append_appends_to_loaded_spec(monkeypatch):
    tool = _tool()
    loaded_spec = {"blocks": [{"type": "heading", "text": "Overview"}]}
    monkeypatch.setattr(
        tool,
        "_load_current",
        lambda raw_id: {"artifact_id": "a-1", "kind": "html", "spec": loaded_spec, "title": "T"},
    )
    captured = {}

    def fake_reversion(artifact_id, kind, spec, action, title=None):
        captured["spec"] = spec
        return {"status": "ok", "artifact_id": artifact_id}

    monkeypatch.setattr(tool, "_reversion", fake_reversion)

    result = tool._edit(id="A1", spec_append={"blocks": [{"type": "paragraph", "text": "Risks…"}]})
    assert result["status"] == "ok"
    assert [b["text"] for b in captured["spec"]["blocks"]] == ["Overview", "Risks…"]


def test_edit_requires_patch_or_append():
    err = _tool()._edit(id="A1")
    assert err["status"] == "error"
    assert "spec_patch and/or spec_append" in err["error"]
