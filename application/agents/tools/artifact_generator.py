"""Artifact Generator tool: render editable documents from a JSON spec and version them append-only.

The ``artifact_versions.spec`` JSONB is the source of truth; the rendered
``.pptx``/``.docx``/``.xlsx``/``.pdf``/``.html`` is derived. ``create_artifact`` stores
v1, ``edit_artifact`` applies an RFC 7386 merge-patch to the current spec and
appends a version, ``rewrite_artifact`` replaces the spec wholesale and appends
a version. Rendering runs a FIXED program in the sandbox that reads the spec as
DATA (``json.loads``) — spec values are never interpolated into the program, so
a spec string containing code/quotes is rendered as literal text, not executed.
"""

from __future__ import annotations

import copy
import json
import logging
import uuid
from typing import Any, Dict, List, Optional

from application.agents.tools.artifact_ref import resolve_artifact_id
from application.agents.tools.base import Tool
from application.core.settings import settings
from application.sandbox.artifacts_capture import (
    QuotaExceeded,
    append_artifact_version,
    persist_new_artifact,
)
from application.sandbox.sandbox_creator import SandboxCreator
from application.storage.db.repositories.artifacts import ArtifactsRepository
from application.storage.db.session import db_readonly

logger = logging.getLogger(__name__)

try:
    import jsonschema
except Exception:  # pragma: no cover - jsonschema is a declared dependency
    jsonschema = None  # type: ignore[assignment]

# Per-kind output metadata: artifact ``kind`` + produced file extension + mime.
_KIND_INFO: Dict[str, Dict[str, str]] = {
    "presentation": {
        "kind": "presentation",
        "ext": "pptx",
        "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    },
    "document": {
        "kind": "document",
        "ext": "docx",
        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    },
    "spreadsheet": {
        "kind": "spreadsheet",
        "ext": "xlsx",
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    },
    "pdf": {
        "kind": "document",
        "ext": "pdf",
        "mime": "application/pdf",
    },
    "html": {
        "kind": "html",
        "ext": "html",
        "mime": "text/html",
    },
}

# Tight per-kind JSON schemas. ``additionalProperties: false`` keeps specs minimal
# and rejects stray keys before any rendering happens.
_SCHEMAS: Dict[str, Dict[str, Any]] = {
    "presentation": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "title": {"type": "string"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "title": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                        "notes": {"type": "string"},
                    },
                },
            },
        },
        "required": ["slides"],
    },
    "document": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "title": {"type": "string"},
            "sections": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "heading": {"type": "string"},
                        "paragraphs": {"type": "array", "items": {"type": "string"}},
                    },
                },
            },
        },
        "required": ["sections"],
    },
    "spreadsheet": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "sheets": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "name": {"type": "string"},
                        "rows": {
                            "type": "array",
                            "items": {"type": "array", "items": {}},
                        },
                    },
                    "required": ["rows"],
                },
            },
        },
        "required": ["sheets"],
    },
    "pdf": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "title": {"type": "string"},
            "blocks": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "type": {"type": "string", "enum": ["heading", "paragraph"]},
                        "text": {"type": "string"},
                    },
                    "required": ["type", "text"],
                },
            },
        },
        "required": ["blocks"],
    },
    "html": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "title": {"type": "string"},
            "blocks": {
                "type": "array",
                "items": {
                    "oneOf": [
                        {
                            "type": "object",
                            "additionalProperties": False,
                            "properties": {
                                "type": {"const": "heading"},
                                "text": {"type": "string"},
                                "level": {"type": "integer", "minimum": 1, "maximum": 3},
                            },
                            "required": ["type", "text"],
                        },
                        {
                            "type": "object",
                            "additionalProperties": False,
                            "properties": {
                                "type": {"const": "paragraph"},
                                "text": {"type": "string"},
                            },
                            "required": ["type", "text"],
                        },
                        {
                            "type": "object",
                            "additionalProperties": False,
                            "properties": {
                                "type": {"const": "list"},
                                "ordered": {"type": "boolean"},
                                "items": {"type": "array", "items": {"type": "string"}},
                            },
                            "required": ["type", "items"],
                        },
                        {
                            "type": "object",
                            "additionalProperties": False,
                            "properties": {
                                "type": {"const": "table"},
                                "headers": {"type": "array", "items": {"type": "string"}},
                                "rows": {
                                    "type": "array",
                                    "items": {"type": "array", "items": {"type": "string"}},
                                },
                            },
                            "required": ["type", "rows"],
                        },
                        {
                            "type": "object",
                            "additionalProperties": False,
                            "properties": {
                                "type": {"const": "code"},
                                "text": {"type": "string"},
                            },
                            "required": ["type", "text"],
                        },
                    ],
                },
            },
        },
        "required": ["blocks"],
    },
}

# FIXED renderer programs. Each reads ``spec.json`` from the workspace as DATA
# and writes ``out.<ext>``. The spec is NEVER string-interpolated into the
# program; ``{spec_path}``/``{out_path}`` are server-controlled path literals.
_RENDERERS: Dict[str, str] = {
    "presentation": (
        "import json\n"
        "from pptx import Presentation\n"
        "from pptx.util import Pt\n"
        "spec = json.load(open({spec_path!r}))\n"
        "prs = Presentation()\n"
        "blank = prs.slide_layouts[6]\n"
        "title_only = prs.slide_layouts[5]\n"
        "for s in spec.get('slides', []):\n"
        "    slide = prs.slides.add_slide(title_only)\n"
        "    slide.shapes.title.text = str(s.get('title', '') or '')\n"
        "    bullets = s.get('bullets') or []\n"
        "    if bullets:\n"
        "        left = top = Pt(72)\n"
        "        width = prs.slide_width - Pt(144)\n"
        "        height = prs.slide_height - Pt(216)\n"
        "        box = slide.shapes.add_textbox(left, Pt(150), width, height)\n"
        "        tf = box.text_frame\n"
        "        tf.word_wrap = True\n"
        "        for i, b in enumerate(bullets):\n"
        "            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()\n"
        "            para.text = str(b)\n"
        "    notes = s.get('notes')\n"
        "    if notes:\n"
        "        slide.notes_slide.notes_text_frame.text = str(notes)\n"
        "prs.save({out_path!r})\n"
    ),
    "document": (
        "import json\n"
        "from docx import Document\n"
        "spec = json.load(open({spec_path!r}))\n"
        "doc = Document()\n"
        "title = spec.get('title')\n"
        "if title:\n"
        "    doc.add_heading(str(title), level=0)\n"
        "for sec in spec.get('sections', []):\n"
        "    heading = sec.get('heading')\n"
        "    if heading:\n"
        "        doc.add_heading(str(heading), level=1)\n"
        "    for p in (sec.get('paragraphs') or []):\n"
        "        doc.add_paragraph(str(p))\n"
        "doc.save({out_path!r})\n"
    ),
    "spreadsheet": (
        "import json\n"
        "from openpyxl import Workbook\n"
        "spec = json.load(open({spec_path!r}))\n"
        "wb = Workbook()\n"
        "wb.remove(wb.active)\n"
        "for idx, sheet in enumerate(spec.get('sheets', [])):\n"
        "    name = str(sheet.get('name') or ('Sheet%d' % (idx + 1)))[:31]\n"
        "    ws = wb.create_sheet(title=name)\n"
        "    for row in (sheet.get('rows') or []):\n"
        "        ws.append([('' if c is None else c) for c in row])\n"
        "if not wb.sheetnames:\n"
        "    wb.create_sheet(title='Sheet1')\n"
        "wb.save({out_path!r})\n"
    ),
    "pdf": (
        "import json\n"
        "from reportlab.lib.pagesizes import letter\n"
        "from reportlab.lib.styles import getSampleStyleSheet\n"
        "from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer\n"
        "from xml.sax.saxutils import escape\n"
        "spec = json.load(open({spec_path!r}))\n"
        "styles = getSampleStyleSheet()\n"
        "story = []\n"
        "title = spec.get('title')\n"
        "if title:\n"
        "    story.append(Paragraph(escape(str(title)), styles['Title']))\n"
        "    story.append(Spacer(1, 12))\n"
        "for block in spec.get('blocks', []):\n"
        "    style = styles['Heading1'] if block.get('type') == 'heading' else styles['BodyText']\n"
        "    story.append(Paragraph(escape(str(block.get('text', ''))), style))\n"
        "    story.append(Spacer(1, 6))\n"
        "SimpleDocTemplate({out_path!r}, pagesize=letter).build(story)\n"
    ),
    "html": (
        "import json\n"
        "import html\n"
        "spec = json.load(open({spec_path!r}))\n"
        "def esc(value):\n"
        "    return html.escape('' if value is None else str(value))\n"
        "parts = []\n"
        "title = spec.get('title')\n"
        "if title:\n"
        "    parts.append('<h1>' + esc(title) + '</h1>')\n"
        "for block in spec.get('blocks', []):\n"
        "    kind = block.get('type')\n"
        "    if kind == 'heading':\n"
        "        level = block.get('level') or 2\n"
        "        try:\n"
        "            level = int(level)\n"
        "        except (TypeError, ValueError):\n"
        "            level = 2\n"
        "        level = min(max(level, 1), 3) + 1\n"
        "        parts.append('<h%d>%s</h%d>' % (level, esc(block.get('text', '')), level))\n"
        "    elif kind == 'paragraph':\n"
        "        parts.append('<p>' + esc(block.get('text', '')) + '</p>')\n"
        "    elif kind == 'list':\n"
        "        tag = 'ol' if block.get('ordered') else 'ul'\n"
        "        items = ''.join('<li>' + esc(i) + '</li>' for i in (block.get('items') or []))\n"
        "        parts.append('<%s>%s</%s>' % (tag, items, tag))\n"
        "    elif kind == 'table':\n"
        "        rows_html = []\n"
        "        headers = block.get('headers')\n"
        "        if headers:\n"
        "            cells = ''.join('<th>' + esc(h) + '</th>' for h in headers)\n"
        "            rows_html.append('<thead><tr>' + cells + '</tr></thead>')\n"
        "        body = []\n"
        "        for row in (block.get('rows') or []):\n"
        "            cells = ''.join('<td>' + esc(c) + '</td>' for c in row)\n"
        "            body.append('<tr>' + cells + '</tr>')\n"
        "        rows_html.append('<tbody>' + ''.join(body) + '</tbody>')\n"
        "        parts.append('<table>' + ''.join(rows_html) + '</table>')\n"
        "    elif kind == 'code':\n"
        "        parts.append('<pre><code>' + esc(block.get('text', '')) + '</code></pre>')\n"
        # CSS braces are doubled so the outer ``_RENDERERS[kind].format(...)`` leaves them literal.
        "css = (\n"
        "    'body{{font-family:system-ui,-apple-system,Segoe UI,Roboto,sans-serif;'\n"
        "    'line-height:1.6;color:#1a1a1a;max-width:800px;margin:0 auto;padding:24px}}'\n"
        "    'h1,h2,h3,h4{{line-height:1.25;margin:1.2em 0 0.5em}}'\n"
        "    'table{{border-collapse:collapse;width:100%;margin:1em 0}}'\n"
        "    'th,td{{border:1px solid #d0d0d0;padding:6px 10px;text-align:left}}'\n"
        "    'th{{background:#f5f5f5}}'\n"
        "    'pre{{background:#f5f5f5;padding:12px;border-radius:6px;overflow:auto}}'\n"
        "    'code{{font-family:ui-monospace,SFMono-Regular,Menlo,monospace}}'\n"
        ")\n"
        "doc = (\n"
        "    '<!doctype html><html lang=\"en\"><head><meta charset=\"utf-8\">'\n"
        "    '<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\">'\n"
        "    '<title>' + esc(title or 'Report') + '</title>'\n"
        "    '<style>' + css + '</style></head><body>'\n"
        "    + ''.join(parts) + '</body></html>'\n"
        ")\n"
        "open({out_path!r}, 'w', encoding='utf-8').write(doc)\n"
    ),
}


def merge_patch(target: Any, patch: Any) -> Any:
    """Apply an RFC 7386 JSON Merge Patch to ``target`` and return the result."""
    if not isinstance(patch, dict):
        return copy.deepcopy(patch)
    if not isinstance(target, dict):
        target = {}
    result = copy.deepcopy(target)
    for key, value in patch.items():
        if value is None:
            result.pop(key, None)
        else:
            result[key] = merge_patch(result.get(key), value)
    return result


class ArtifactGeneratorTool(Tool):
    """Artifact
    Create, edit, and version documents - slides, docs, sheets, PDF, HTML.
    """

    def __init__(self, tool_config: Optional[Dict[str, Any]] = None, user_id: Optional[str] = None) -> None:
        """Bind the tool to the invoker and its conversation/run-scoped sandbox session."""
        self.config: Dict[str, Any] = tool_config or {}
        self.user_id: Optional[str] = user_id
        self.tool_id: Optional[str] = self.config.get("tool_id")
        self.conversation_id: Optional[str] = self.config.get("conversation_id")
        self.workflow_run_id: Optional[str] = self.config.get("workflow_run_id")
        self._last_artifact_id: Optional[str] = None

    # ------------------------------------------------------------------
    # Tool ABC
    # ------------------------------------------------------------------
    def get_actions_metadata(self) -> List[Dict[str, Any]]:
        """Return JSON metadata describing the create/edit/rewrite actions for tool schemas."""
        kinds = sorted(_KIND_INFO.keys())
        return [
            {
                "name": "create_artifact",
                "description": (
                    "Render a new editable document from a JSON spec and store it as version 1. "
                    "The spec is the source of truth; the rendered file is derived. The response "
                    "carries a short ref (like `A1`) you can pass to edit_artifact/rewrite_artifact."
                ),
                "active": True,
                "parameters": {
                    "type": "object",
                    "properties": {
                        "kind": {
                            "type": "string",
                            "enum": kinds,
                            "description": (
                                "Document kind to render; `html` is an inline-rendered, versionable HTML report."
                            ),
                        },
                        "title": {"type": "string", "description": "Optional artifact title."},
                        "spec": {"type": "object", "description": "Document spec matching the kind's schema."},
                    },
                    "required": ["kind", "spec"],
                },
            },
            {
                "name": "edit_artifact",
                "description": (
                    "Apply a JSON merge-patch (RFC 7386) to the current spec, re-render, and append a "
                    "new version. Preferred for small, targeted changes."
                ),
                "active": True,
                "parameters": {
                    "type": "object",
                    "properties": {
                        "id": {
                            "type": "string",
                            "description": "Artifact to edit; accepts the short ref like `A1` "
                            "returned by a previous artifact action, or the full artifact id.",
                        },
                        "spec_patch": {
                            "type": "object",
                            "description": "RFC 7386 merge-patch; null values delete keys.",
                        },
                    },
                    "required": ["id", "spec_patch"],
                },
            },
            {
                "name": "rewrite_artifact",
                "description": "Replace the spec wholesale, re-render, and append a new version.",
                "active": True,
                "parameters": {
                    "type": "object",
                    "properties": {
                        "id": {
                            "type": "string",
                            "description": "Artifact to rewrite; accepts the short ref like `A1` "
                            "returned by a previous artifact action, or the full artifact id.",
                        },
                        "spec": {"type": "object", "description": "Replacement spec matching the kind's schema."},
                    },
                    "required": ["id", "spec"],
                },
            },
        ]

    def get_config_requirements(self) -> Dict[str, Any]:
        """Return configuration requirements (none beyond the deployment sandbox backend)."""
        return {}

    def get_artifact_id(self, action_name: str, **kwargs: Any) -> Optional[str]:
        """Return the produced artifact id so the UI artifact rail lights up."""
        return self._last_artifact_id

    # ------------------------------------------------------------------
    # Dispatch
    # ------------------------------------------------------------------
    def execute_action(self, action_name: str, **kwargs: Any) -> Dict[str, Any]:
        """Dispatch a create/edit/rewrite action."""
        self._last_artifact_id = None
        if not self.user_id:
            return {"status": "error", "error": "artifact_generator requires a valid user_id."}
        if self.conversation_id is None and self.workflow_run_id is None:
            return {"status": "error", "error": "artifact_generator requires a conversation_id or workflow_run_id."}
        if jsonschema is None:
            return {"status": "error", "error": "jsonschema is required for spec validation."}
        if action_name == "create_artifact":
            return self._create(**kwargs)
        if action_name == "edit_artifact":
            return self._edit(**kwargs)
        if action_name == "rewrite_artifact":
            return self._rewrite(**kwargs)
        return {"status": "error", "error": f"unknown action: {action_name}"}

    # ------------------------------------------------------------------
    # Actions
    # ------------------------------------------------------------------
    def _create(self, **kwargs: Any) -> Dict[str, Any]:
        """Validate, render, and persist a new artifact at version 1."""
        kind = kwargs.get("kind")
        spec = kwargs.get("spec")
        title = kwargs.get("title")
        if kind not in _KIND_INFO:
            return {"status": "error", "error": f"unsupported kind: {kind!r}; expected one of {sorted(_KIND_INFO)}."}
        valid = self._validate(kind, spec)
        if valid is not None:
            return valid

        rendered = self._render(kind, spec)
        if rendered.get("error"):
            return {"status": "error", "error": rendered["error"]}

        info = _KIND_INFO[kind]
        filename = self._filename(title, kind)
        try:
            ref = persist_new_artifact(
                user_id=self.user_id,
                kind=info["kind"],
                data=rendered["data"],
                filename=filename,
                mime_type=info["mime"],
                title=title,
                conversation_id=self.conversation_id,
                workflow_run_id=self.workflow_run_id,
                spec=spec,
                produced_by=self._produced_by("create_artifact", kind),
            )
        except QuotaExceeded as exc:
            return {"status": "error", "error": str(exc)}
        if ref is None:
            return {"status": "error", "error": "failed to persist artifact."}
        self._last_artifact_id = ref["artifact_id"]
        return {"status": "ok", **ref}

    def _edit(self, **kwargs: Any) -> Dict[str, Any]:
        """Merge-patch the current spec, re-render, and append a version."""
        spec_patch = kwargs.get("spec_patch")
        if not isinstance(spec_patch, dict):
            return {"status": "error", "error": "spec_patch must be a JSON object (merge-patch)."}
        loaded = self._load_current(kwargs.get("id"))
        if loaded.get("error"):
            return {"status": "error", "error": loaded["error"]}
        new_spec = merge_patch(loaded["spec"], spec_patch)
        return self._reversion(loaded["artifact_id"], loaded["kind"], new_spec, "edit_artifact")

    def _rewrite(self, **kwargs: Any) -> Dict[str, Any]:
        """Replace the spec wholesale, re-render, and append a version."""
        spec = kwargs.get("spec")
        loaded = self._load_current(kwargs.get("id"))
        if loaded.get("error"):
            return {"status": "error", "error": loaded["error"]}
        return self._reversion(loaded["artifact_id"], loaded["kind"], spec, "rewrite_artifact")

    def _reversion(self, artifact_id: str, kind: str, spec: Any, action: str) -> Dict[str, Any]:
        """Validate the new spec, re-render, and append the next version of an existing artifact."""
        valid = self._validate(kind, spec)
        if valid is not None:
            return valid
        rendered = self._render(kind, spec)
        if rendered.get("error"):
            return {"status": "error", "error": rendered["error"]}
        info = _KIND_INFO[kind]
        filename = self._filename(None, kind)
        try:
            ref = append_artifact_version(
                user_id=self.user_id,
                artifact_id=artifact_id,
                data=rendered["data"],
                filename=filename,
                mime_type=info["mime"],
                spec=spec,
                produced_by=self._produced_by(action, kind),
                conversation_id=self.conversation_id,
                workflow_run_id=self.workflow_run_id,
            )
        except QuotaExceeded as exc:
            return {"status": "error", "error": str(exc)}
        if ref is None:
            return {"status": "error", "error": "failed to persist artifact version."}
        self._last_artifact_id = ref["artifact_id"]
        return {"status": "ok", **ref}

    # ------------------------------------------------------------------
    # Spec / render helpers
    # ------------------------------------------------------------------
    def _validate(self, kind: str, spec: Any) -> Optional[Dict[str, Any]]:
        """Return an error payload when ``spec`` is invalid for ``kind``, else None."""
        if not isinstance(spec, dict):
            return {"status": "error", "error": "spec must be a JSON object."}
        try:
            jsonschema.validate(spec, _SCHEMAS[kind])
        except jsonschema.ValidationError as exc:
            return {"status": "error", "error": f"invalid {kind} spec: {exc.message}"}
        return None

    def _load_current(self, raw_id: Any) -> Dict[str, Any]:
        """Resolve a short ref/uuid to its parent-scoped artifact and current-version spec for edit/rewrite."""
        if not isinstance(raw_id, str) or not raw_id.strip():
            return {"error": "id is required."}
        try:
            with db_readonly() as conn:
                repo = ArtifactsRepository(conn)
                # A ref (A1/A2/...) resolves to an id within this parent only; the
                # resolved id is then re-checked through the parent-scoped gate so a
                # ref can never reach another tenant.
                artifact_id = resolve_artifact_id(
                    repo,
                    raw_id.strip(),
                    conversation_id=self.conversation_id,
                    workflow_run_id=self.workflow_run_id,
                )
                if artifact_id is None:
                    return {"error": f"artifact {raw_id} not found in this conversation/run."}
                artifact = repo.get_artifact_in_parent(
                    artifact_id,
                    conversation_id=self.conversation_id,
                    workflow_run_id=self.workflow_run_id,
                )
                if artifact is None:
                    return {"error": f"artifact {raw_id} not found in this conversation/run."}
                version = repo.get_version(artifact_id, artifact["current_version"])
        except Exception:
            logger.exception("artifact_generator: failed to load artifact")
            return {"error": f"failed to load artifact {raw_id}."}
        if not version or version.get("spec") is None:
            return {"error": f"artifact {raw_id} has no editable spec."}
        kind = self._kind_for(artifact, version)
        if kind is None:
            return {"error": f"artifact {raw_id} is not a spec-rendered document."}
        return {"artifact_id": artifact_id, "kind": kind, "spec": version["spec"]}

    @staticmethod
    def _kind_for(artifact: Dict[str, Any], version: Dict[str, Any]) -> Optional[str]:
        """Resolve the spec kind from ``produced_by`` (preferred) or the version mime type."""
        produced = version.get("produced_by")
        if isinstance(produced, dict):
            spec_kind = produced.get("spec_kind")
            if spec_kind in _KIND_INFO:
                return spec_kind
        mime = version.get("mime_type") or ""
        for spec_kind, info in _KIND_INFO.items():
            if info["mime"] == mime:
                return spec_kind
        return None

    def _render(self, kind: str, spec: Any) -> Dict[str, Any]:
        """Run the fixed renderer in the sandbox and return the produced file bytes."""
        session_id = self._resolve_session_id()
        if session_id is None:
            return {"error": "artifact_generator requires a conversation_id or workflow_run_id."}

        token = uuid.uuid4().hex
        token_dir = f"artifacts/{token}"
        spec_path = f"{token_dir}/spec.json"
        out_path = f"{token_dir}/out.{_KIND_INFO[kind]['ext']}"
        program = _RENDERERS[kind].format(spec_path=spec_path, out_path=out_path)
        timeout = float(getattr(settings, "SANDBOX_EXEC_TIMEOUT", 60))

        manager = SandboxCreator.get_manager()
        try:
            manager.open(session_id, ttl=timeout)
        except Exception as exc:
            logger.exception("artifact_generator: failed to open sandbox session")
            return {"error": f"sandbox unavailable: {type(exc).__name__}: {exc}"}
        try:
            # The spec rides in as a JSON file the program ``json.load``s; it is
            # never interpolated into the program, so its contents stay data.
            manager.put_file(session_id, spec_path, json.dumps(spec).encode("utf-8"))
            result = manager.exec(session_id, program, timeout=timeout)
            if not result.ok:
                detail = (
                    f"{result.error_name}: {result.error_value}"
                    if result.error_name
                    else (result.error_value or "render failed")
                )
                return {"error": f"render failed: {detail}"}
            data = manager.get_file(session_id, out_path)
        except Exception as exc:
            logger.exception("artifact_generator: render failed")
            return {"error": f"render failed: {type(exc).__name__}: {exc}"}
        finally:
            # Drop this render's scratch dir before tearing the session down so a
            # warm/reused session doesn't accumulate per-render files on disk.
            manager.remove_path(session_id, token_dir)
            try:
                manager.close(session_id)
            except Exception:
                logger.exception("artifact_generator: session close failed")
        if not data:
            return {"error": "renderer produced an empty file."}
        return {"data": data}

    # ------------------------------------------------------------------
    # Misc helpers
    # ------------------------------------------------------------------
    def _produced_by(self, action: str, kind: str) -> Dict[str, Any]:
        """Build the ``produced_by`` provenance record, carrying the spec kind for re-editing."""
        return {
            "tool": "artifact_generator",
            "action": action,
            "spec_kind": kind,
            "tool_id": self.tool_id,
        }

    @staticmethod
    def _filename(title: Optional[str], kind: str) -> str:
        """Derive a download filename from a title (or a generic stem) plus the kind extension."""
        if kind == "html":
            return "report.html"
        stem = (title or "artifact").strip() or "artifact"
        return f"{stem}.{_KIND_INFO[kind]['ext']}"

    def _resolve_session_id(self) -> Optional[str]:
        """Derive the sandbox session id from the bound conversation/run; sanitize to the gateway charset."""
        raw = self.conversation_id or self.workflow_run_id
        if not raw:
            return None
        sanitized = "".join(c if c.isalnum() or c in "-_" else "-" for c in str(raw))
        return sanitized or None
