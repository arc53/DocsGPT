"""PPT parser.
Contains parsers for presentation (.pptx) files to extract slide text.
"""
from pathlib import Path
from typing import Any, Dict, List, Union

from application.parser.file.base_parser import BaseParser

class PPTXParser(BaseParser):
    r"""PPTX (.pptx) parser for extracting text from PowerPoint slides.
    Args:
        concat_slides (bool): Specifies whether to concatenate all slide text into one document.
            - If True, slide texts will be joined together as a single string.
            - If False, each slide's text will be stored as a separate entry in a list.
            Set to True by default.
        slide_separator (str): Separator used to join slides' text content.
            Only used when `concat_slides=True`. Default is "\n".
        Refer to https://python-pptx.readthedocs.io/en/latest/ for more information.
    """

    def __init__(
        self,
        *args: Any,
        concat_slides: bool = True,
        slide_separator: str = "\n",
        **kwargs: Any
    ) -> None:
        """Init params."""
        super().__init__(*args, **kwargs)
        self._concat_slides = concat_slides
        self._slide_separator = slide_separator

    def _init_parser(self) -> Dict:
        """Init parser."""
        return {}

    def parse_file(self, file: Path, errors: str = "ignore") -> Union[str, List[str]]:
        r"""
        Parse a .pptx file and extract text from each slide.
        Args:
            file (Path): Path to the .pptx file.
            errors (str): Error handling policy ('ignore' by default).
        Returns:
            Union[str, List[str]]: Concatenated text if concat_slides is True,
            otherwise a list of slide texts.
        """

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("pptx module is required to read .PPTX files.")

        try:
            presentation = Presentation(file)
            slide_texts=[]

            # Iterate over each slide in the presentation
            for slide in presentation.slides:
                slide_text=""

                # Iterate over each shape in the slide
                for shape in slide.shapes:
                    # Check if the shape has a 'text' attribute and append that to the slide_text
                    if hasattr(shape,"text"):
                        slide_text+=shape.text

                slide_texts.append(slide_text.strip())

            if self._concat_slides:
                return self._slide_separator.join(slide_texts)
            else:
                return slide_texts

        except Exception as e:
            raise e